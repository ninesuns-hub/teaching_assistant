import os
from pptx import Presentation
import sys

# Set stdout to utf-8 to avoid encoding errors in some environments
if sys.stdout.encoding != 'utf-8':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_pptx_text(file_path):
    if not os.path.exists(file_path):
        return f"File not found: {file_path}"

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return f"Error loading presentation: {e}"

    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_runs.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

    return "\n\n".join(text_runs)

if __name__ == "__main__":
    if len(sys.argv) > 1:
        print(extract_pptx_text(sys.argv[1]))
