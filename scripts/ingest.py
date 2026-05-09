import sys
import os
import argparse

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import fitz
from pptx import Presentation
from config import CHUNK_SIZE, CHUNK_OVERLAP, DATA_DIR
from database.vector_store import add_documents, get_stats


def split_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    """按字数切片，保留重叠避免语义断裂。"""
    if not text.strip():
        return []
    chunks = []
    start = 0
    text = text.strip()
    while start < len(text):
        chunk = text[start:start + chunk_size].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


def parse_pptx(file_path: str, chapter: str) -> list[dict]:
    """每页 PPT 作为一个独立片段，保留完整页面语义。"""
    prs = Presentation(file_path)
    filename = os.path.basename(file_path)
    chunks = []

    for slide_idx, slide in enumerate(prs.slides):
        texts = []

        for shape in slide.shapes:
            # ── 提取文本框内容 ──
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            # ── 提取表格内容──
            if shape.has_table:
                tbl = shape.table
                num_rows = len(tbl.rows)
                if num_rows == 0:
                    continue

                # 提取表头
                headers = [
                    cell.text.strip().replace('\n', ' ')
                    for cell in tbl.rows[0].cells
                ]

                # 从第1行开始遍历数据行（用索引代替切片）
                for row_idx in range(1, num_rows):
                    row = tbl.rows[row_idx]
                    row_parts = []
                    for header, cell in zip(headers, row.cells):
                        cell_text = cell.text.strip().replace('\n', ' ')
                        if cell_text:
                            row_parts.append(f"{header}: {cell_text}")
                    if row_parts:
                        texts.append(" | ".join(row_parts))

        page_text = "\n".join(texts).strip()
        if not page_text:
            continue

        chunks.append({
            "id":          f"{chapter}_ppt_{slide_idx}",
            "text":        page_text,
            "source_type": "ppt",
            "source_file": filename,
            "chapter":     chapter,
            "page":        slide_idx + 1,
        })

    print(f"[PPT] {filename} 共提取 {len(chunks)} 页片段")
    return chunks


def parse_pdf(file_path: str, chapter: str) -> list[dict]:
    """
    按页提取，每页内容如果太长再切片。
    """
    doc = fitz.open(file_path)
    filename = os.path.basename(file_path)
    chunks = []
    chunk_global_idx = 0  # 全局片段编号，保证 id 唯一

    for page_num, page in enumerate(doc):
        page_text = page.get_text("text").strip()
        if not page_text:
            continue

        real_page = page_num + 1  # 真实页码从1开始

        # 如果这一页内容不超过 CHUNK_SIZE，整页作为一个片段
        if len(page_text) <= CHUNK_SIZE:
            chunks.append({
                "id":          f"{chapter}_pdf_{chunk_global_idx}",
                "text":        page_text,
                "source_type": "pdf",
                "source_file": filename,
                "chapter":     chapter,
                "page":        real_page,
            })
            chunk_global_idx += 1
        else:
            # 页面内容较长时再切片，每个子片段都记录真实页码
            sub_chunks = split_text(page_text, CHUNK_SIZE, CHUNK_OVERLAP)
            for sub in sub_chunks:
                chunks.append({
                    "id":          f"{chapter}_pdf_{chunk_global_idx}",
                    "text":        sub,
                    "source_type": "pdf",
                    "source_file": filename,
                    "chapter":     chapter,
                    "page":        real_page,  # 同一页的子片段共享真实页码
                })
                chunk_global_idx += 1

    doc.close()
    print(f"[PDF] {filename} 共切出 {len(chunks)} 个片段（按真实页码）")
    return chunks


def ingest_file(file_path: str, chapter: str) -> None:
    if not os.path.exists(file_path):
        print(f"[错误] 文件不存在：{file_path}")
        return

    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        chunks = parse_pptx(file_path, chapter)
    elif ext == ".pdf":
        chunks = parse_pdf(file_path, chapter)
    else:
        print(f"[跳过] 不支持的文件类型：{ext}")
        return

    if chunks:
        add_documents(chunks)


def ingest_all(docs_dir: str = "docs") -> None:
    if not os.path.exists(docs_dir):
        print(f"[错误] 目录不存在：{docs_dir}")
        return

    for filename in sorted(os.listdir(docs_dir)):
        # 跳过 Office 临时文件
        if filename.startswith("~$"):
            continue
        
        ext = os.path.splitext(filename)[1].lower()
        if ext not in (".pptx", ".pdf"):
            continue
        file_path = os.path.join(docs_dir, filename)
        chapter = filename.split("-")[0].split("_")[0]
        print(f"\n正在导入：{filename}（章节：{chapter}）")
        ingest_file(file_path, chapter)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="向向量数据库导入课程文件")
    parser.add_argument("--file",    help="单个文件路径")
    parser.add_argument("--chapter", help="章节名，如 第0章、第1章")
    parser.add_argument("--all",     action="store_true", help="导入 docs/ 目录下所有文件")
    parser.add_argument("--stats",   action="store_true", help="查看数据库状态")
    args = parser.parse_args()

    os.makedirs(DATA_DIR, exist_ok=True)

    if args.stats:
        stats = get_stats()
        print(f"\n数据库状态：总片段数 {stats['total']}（PPT: {stats['ppt_count']}，PDF: {stats['pdf_count']}）")
    elif args.all:
        ingest_all()
        stats = get_stats()
        print(f"\n导入完成：总片段数 {stats['total']}（PPT: {stats['ppt_count']}，PDF: {stats['pdf_count']}）")
    elif args.file:
        chapter = args.chapter or "未知章节"
        ingest_file(args.file, chapter)
        stats = get_stats()
        print(f"\n导入完成：总片段数 {stats['total']}（PPT: {stats['ppt_count']}，PDF: {stats['pdf_count']}）")
    else:
        parser.print_help()