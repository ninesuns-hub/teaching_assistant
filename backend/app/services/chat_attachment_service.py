from __future__ import annotations

import base64
import html
import logging
import os
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from pathlib import Path
from typing import Callable

import fitz
from docx import Document
from openai import OpenAI
from pptx import Presentation
from sqlalchemy import inspect

from agent_core.config.settings import settings
from database import chat_attachment_repo
from database.mysql_db import ChatMessageAttachment, SessionLocal


logger = logging.getLogger(__name__)

ALLOWED_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppsx"}
DOCUMENT_MIME_TYPES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppsx": "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
}
MAX_DOCUMENT_FILE_SIZE = 20 * 1024 * 1024
MAX_MESSAGE_TOTAL_SIZE = 50 * 1024 * 1024
MAX_MESSAGE_ATTACHMENTS = 3
MAX_OFFICE_UNCOMPRESSED_SIZE = 100 * 1024 * 1024
MAX_OFFICE_ARCHIVE_ENTRIES = 5000

_executor = ThreadPoolExecutor(max_workers=1, thread_name_prefix="chat-attachment")


class AttachmentCancelled(Exception):
    pass


class AttachmentParseError(ValueError):
    pass


def safe_filename(value: str) -> str:
    name = Path(value or "document").name.strip()
    return (name or "document")[:255]


def extension_for(filename: str) -> str:
    return Path(filename).suffix.lower()


def validate_document_bytes(filename: str, raw: bytes) -> tuple[str, str]:
    ext = extension_for(filename)
    if ext in {".doc", ".ppt"}:
        raise AttachmentParseError("旧版 .doc/.ppt 暂不支持，请另存为 .docx/.pptx 后上传")
    if ext not in ALLOWED_DOCUMENT_EXTENSIONS:
        raise AttachmentParseError("仅支持 PDF、DOCX、PPTX 和 PPSX 文档")
    if not raw:
        raise AttachmentParseError("文档内容为空")
    if len(raw) > MAX_DOCUMENT_FILE_SIZE:
        raise AttachmentParseError("单个文档不能超过 20MB")
    if ext == ".pdf" and not raw.startswith(b"%PDF"):
        raise AttachmentParseError("文件内容不是有效的 PDF")
    if ext != ".pdf":
        if not raw.startswith(b"PK"):
            raise AttachmentParseError("Office 文件格式无效或文件已损坏")
        try:
            with zipfile.ZipFile(BytesIO(raw)) as archive:
                entries = archive.infolist()
                if len(entries) > MAX_OFFICE_ARCHIVE_ENTRIES:
                    raise AttachmentParseError("Office 文件内部条目过多")
                if any(item.flag_bits & 0x1 for item in entries):
                    raise AttachmentParseError("加密 Office 文件暂不支持")
                if sum(item.file_size for item in entries) > MAX_OFFICE_UNCOMPRESSED_SIZE:
                    raise AttachmentParseError("Office 文件解压后体积过大")
                if archive.testzip() is not None:
                    raise AttachmentParseError("Office 文件已损坏")
        except AttachmentParseError:
            raise
        except (OSError, zipfile.BadZipFile) as exc:
            raise AttachmentParseError("Office 文件格式无效或文件已损坏") from exc
    return ext, DOCUMENT_MIME_TYPES[ext]


def _attachment_root() -> str:
    root = os.path.abspath(settings.CHAT_ATTACHMENTS_DIR)
    os.makedirs(root, exist_ok=True)
    return root


def save_attachment_file(user_id: int, public_id: str, ext: str, raw: bytes) -> str:
    relative_path = f"{user_id}/{public_id}{ext}"
    absolute_path = resolve_attachment_path(relative_path)
    os.makedirs(os.path.dirname(absolute_path), exist_ok=True)
    with open(absolute_path, "wb") as output:
        output.write(raw)
    return relative_path


def resolve_attachment_path(relative_path: str) -> str:
    root = _attachment_root()
    normalized = str(relative_path or "").replace("\\", "/").lstrip("/")
    candidate = os.path.abspath(os.path.join(root, normalized))
    try:
        contained = os.path.commonpath([root, candidate]) == root
    except ValueError:
        contained = False
    if not contained:
        raise AttachmentParseError("附件路径无效")
    return candidate


def delete_attachment_file(relative_path: str) -> None:
    try:
        path = resolve_attachment_path(relative_path)
        if os.path.isfile(path):
            os.remove(path)
    except Exception:
        logger.exception("删除会话附件失败 path=%s", relative_path)


def serialize_attachment(attachment: ChatMessageAttachment) -> dict:
    return {
        "id": attachment.public_id,
        "filename": attachment.filename,
        "file_type": attachment.file_type,
        "mime_type": attachment.mime_type,
        "file_size": int(attachment.file_size or 0),
        "status": attachment.status,
        "progress_current": int(attachment.progress_current or 0),
        "progress_total": int(attachment.progress_total or 0),
        "truncated": bool(attachment.extraction_truncated),
        "requires_ocr": bool(attachment.requires_ocr),
        "error_message": attachment.error_message,
        "created_at": attachment.created_at.isoformat() if attachment.created_at else None,
    }


def submit_attachment_job(public_id: str) -> None:
    _executor.submit(_run_attachment_job, public_id)


def recover_incomplete_chat_attachments() -> int:
    db = SessionLocal()
    try:
        if not inspect(db.get_bind()).has_table("chat_message_attachments"):
            logger.info("会话附件表尚未创建，跳过解析任务恢复")
            return 0
        public_ids = chat_attachment_repo.recover_incomplete(db)
    except Exception:
        db.rollback()
        logger.exception("恢复会话附件解析任务失败")
        return 0
    finally:
        db.close()
    for public_id in public_ids:
        submit_attachment_job(public_id)
    if public_ids:
        logger.info("已恢复会话附件解析任务 count=%s", len(public_ids))
    return len(public_ids)


def cleanup_expired_chat_attachments() -> int:
    db = SessionLocal()
    cleaned = 0
    try:
        if not inspect(db.get_bind()).has_table("chat_message_attachments"):
            return 0
        expired = chat_attachment_repo.expired_unassociated(
            db,
            retention_hours=settings.CHAT_ATTACHMENT_RETENTION_HOURS,
        )
        for attachment in expired:
            delete_attachment_file(attachment.file_path)
            db.delete(attachment)
            cleaned += 1
        db.commit()
    except Exception:
        db.rollback()
        logger.exception("清理过期会话附件失败")
    finally:
        db.close()
    return cleaned


def _check_cancelled(db, public_id: str) -> None:
    if chat_attachment_repo.is_cancelled(db, public_id):
        raise AttachmentCancelled()


def _append_chunk(
    chunks: list[dict],
    *,
    location: str,
    text: str,
    current_chars: int,
) -> tuple[int, bool]:
    normalized = "\n".join(
        line.rstrip()
        for line in str(text or "").replace("\x00", "").splitlines()
        if line.strip()
    ).strip()
    if not normalized:
        return current_chars, False
    remaining = settings.CHAT_ATTACHMENT_MAX_EXTRACTED_CHARS - current_chars
    if remaining <= 0:
        return current_chars, True
    truncated = len(normalized) > remaining
    content = normalized[:remaining]
    chunks.append({"location": location, "text": content})
    return current_chars + len(content), truncated


def _ocr_pdf_page(page, filename: str, page_number: int) -> str:
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
    raw = pixmap.tobytes("jpeg", jpg_quality=82)
    encoded = base64.b64encode(raw).decode("ascii")
    prompt = (
        "你是文档OCR助手。请忠实转写这一页中的全部可见文字、数学公式、表格标题和图注；"
        "保留自然阅读顺序，数学表达尽量使用LaTeX。不要回答或解释文档内容，"
        f"只输出转写结果。文件：{filename}，第{page_number}页。"
    )
    client = OpenAI(
        api_key=settings.VISION_API_KEY or settings.EMBED_API_KEY,
        base_url=settings.VISION_BASE_URL or settings.EMBED_BASE_URL,
        timeout=90.0,
    )
    last_error: Exception | None = None
    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model=settings.VISION_MODEL_NAME,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{encoded}",
                                "detail": "auto",
                            },
                        },
                    ],
                }],
                max_tokens=2000,
                temperature=0,
            )
            text = (response.choices[0].message.content or "").strip()
            if text:
                return text
            raise RuntimeError("OCR 返回内容为空")
        except Exception as exc:
            last_error = exc
            if attempt < 2:
                time.sleep(attempt + 1)
    raise AttachmentParseError(f"第 {page_number} 页 OCR 失败：{last_error}")


def _extract_pdf(
    db,
    attachment: ChatMessageAttachment,
    path: str,
) -> tuple[list[dict], bool]:
    try:
        document = fitz.open(path)
    except Exception as exc:
        raise AttachmentParseError(f"PDF 无法打开：{exc}") from exc
    try:
        if document.needs_pass:
            raise AttachmentParseError("加密 PDF 暂不支持")
        total_pages = document.page_count
        if total_pages <= 0:
            raise AttachmentParseError("PDF 没有可读取页面")
        chunks: list[dict] = []
        char_count = 0
        truncated = False
        ocr_initialized = False
        for index in range(total_pages):
            _check_cancelled(db, attachment.public_id)
            page = document.load_page(index)
            text = (page.get_text("text") or "").strip()
            needs_ocr = len("".join(text.split())) < 20
            if needs_ocr:
                if total_pages > settings.CHAT_ATTACHMENT_OCR_MAX_PAGES:
                    raise AttachmentParseError(
                        f"扫描 PDF 最多支持 {settings.CHAT_ATTACHMENT_OCR_MAX_PAGES} 页，请拆分后上传"
                    )
                if not ocr_initialized:
                    if chat_attachment_repo.count_recent_ocr_documents(
                        db, attachment.user_id
                    ) >= settings.CHAT_ATTACHMENT_OCR_PER_HOUR:
                        raise AttachmentParseError("每小时最多解析 3 份扫描 PDF，请稍后再试")
                    chat_attachment_repo.update_progress(
                        db,
                        attachment.public_id,
                        current=index,
                        total=total_pages,
                        requires_ocr=True,
                    )
                    ocr_initialized = True
                text = _ocr_pdf_page(
                    page,
                    attachment.filename,
                    index + 1,
                )
            char_count, hit_limit = _append_chunk(
                chunks,
                location=f"第 {index + 1} 页",
                text=text,
                current_chars=char_count,
            )
            truncated = truncated or hit_limit
            chat_attachment_repo.update_progress(
                db,
                attachment.public_id,
                current=index + 1,
                total=total_pages,
            )
            if hit_limit:
                break
        if not chunks:
            raise AttachmentParseError("PDF 中未识别到可用文字")
        return chunks, truncated
    finally:
        document.close()


def _extract_docx(
    db,
    attachment: ChatMessageAttachment,
    path: str,
) -> tuple[list[dict], bool]:
    try:
        document = Document(path)
    except Exception as exc:
        raise AttachmentParseError(f"DOCX 无法打开：{exc}") from exc
    units: list[tuple[str, str]] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        if paragraph.text.strip():
            units.append((f"段落 {index}", paragraph.text))
    for table_index, table in enumerate(document.tables, start=1):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in table.rows
        ]
        if any(row.strip(" |") for row in rows):
            units.append((f"表格 {table_index}", "\n".join(rows)))
    if not units:
        raise AttachmentParseError("DOCX 中未提取到文字；图片型文档请转换为扫描 PDF")
    chunks: list[dict] = []
    char_count = 0
    truncated = False
    total = len(units)
    chat_attachment_repo.update_progress(
        db, attachment.public_id, current=0, total=total
    )
    for index, (location, text) in enumerate(units, start=1):
        _check_cancelled(db, attachment.public_id)
        char_count, hit_limit = _append_chunk(
            chunks,
            location=location,
            text=text,
            current_chars=char_count,
        )
        truncated = truncated or hit_limit
        chat_attachment_repo.update_progress(
            db, attachment.public_id, current=index, total=total
        )
        if hit_limit:
            break
    return chunks, truncated


def _extract_presentation(
    db,
    attachment: ChatMessageAttachment,
    path: str,
) -> tuple[list[dict], bool]:
    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise AttachmentParseError(f"演示文稿无法打开：{exc}") from exc
    total = len(presentation.slides)
    if total <= 0:
        raise AttachmentParseError("演示文稿没有幻灯片")
    chunks: list[dict] = []
    char_count = 0
    truncated = False
    chat_attachment_repo.update_progress(
        db, attachment.public_id, current=0, total=total
    )
    for index, slide in enumerate(presentation.slides, start=1):
        _check_cancelled(db, attachment.public_id)
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        char_count, hit_limit = _append_chunk(
            chunks,
            location=f"第 {index} 张幻灯片",
            text="\n".join(texts),
            current_chars=char_count,
        )
        truncated = truncated or hit_limit
        chat_attachment_repo.update_progress(
            db, attachment.public_id, current=index, total=total
        )
        if hit_limit:
            break
    if not chunks:
        raise AttachmentParseError("演示文稿中未提取到文字")
    return chunks, truncated


def _run_attachment_job(public_id: str) -> None:
    db = SessionLocal()
    attachment: ChatMessageAttachment | None = None
    try:
        if not chat_attachment_repo.claim_attachment(db, public_id):
            return
        attachment = chat_attachment_repo.get_attachment(db, public_id)
        if not attachment:
            return
        path = resolve_attachment_path(attachment.file_path)
        extractor: Callable = {
            "pdf": _extract_pdf,
            "docx": _extract_docx,
            "pptx": _extract_presentation,
            "ppsx": _extract_presentation,
        }.get(attachment.file_type)
        if extractor is None:
            raise AttachmentParseError("不支持的文档类型")
        chunks, truncated = extractor(db, attachment, path)
        _check_cancelled(db, public_id)
        chat_attachment_repo.complete_attachment(
            db,
            public_id,
            extracted_content=chunks,
            truncated=truncated,
        )
    except AttachmentCancelled:
        if attachment:
            delete_attachment_file(attachment.file_path)
        chat_attachment_repo.delete_attachment_record(db, public_id)
    except Exception as exc:
        logger.exception("会话附件解析失败 attachment_id=%s", public_id)
        chat_attachment_repo.fail_attachment(
            db,
            public_id,
            str(exc) or "文档解析失败",
        )
    finally:
        db.close()


def office_preview_html(attachment: ChatMessageAttachment) -> str:
    title = html.escape(attachment.filename)
    sections = []
    for item in attachment.extracted_content or []:
        location = html.escape(str(item.get("location") or ""))
        text = html.escape(str(item.get("text") or ""))
        sections.append(
            f"<section><h2>{location}</h2><pre>{text}</pre></section>"
        )
    body = "".join(sections) or "<p>没有可显示的文字内容。</p>"
    truncated = (
        "<p class='notice'>文档内容较长，预览与问答上下文已截断。</p>"
        if attachment.extraction_truncated
        else ""
    )
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width">
<title>{title}</title>
<style>
body{{margin:0;padding:28px;font-family:system-ui,sans-serif;color:#172b3c;background:#f7fafc}}
main{{max-width:900px;margin:auto}}h1{{font-size:22px;overflow-wrap:anywhere}}
section{{margin:16px 0;padding:18px;border:1px solid #d7e0e8;border-radius:12px;background:white}}
h2{{font-size:14px;color:#527086}}pre{{white-space:pre-wrap;overflow-wrap:anywhere;font:15px/1.7 system-ui,sans-serif}}
.notice{{padding:12px;border-radius:10px;background:#fff3cd;color:#6a5312}}
</style></head><body><main><h1>{title}</h1>{truncated}{body}</main></body></html>"""
