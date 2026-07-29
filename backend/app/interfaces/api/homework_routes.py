import json
import mimetypes
import os
import shutil
from datetime import datetime
from html import escape
from typing import Optional

from docx import Document
from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from sqlalchemy.orm import Session

from app.core.data_manager import data_manager
from app.core.deps import get_current_user, require_role
from agent_core.config.settings import settings
from database import class_repo, homework_repo
from database.mysql_db import User, UserRole, get_db
from .schemas import HomeworkAttachmentResponse, HomeworkResponse, HomeworkSubmissionResponse

router = APIRouter()

ALLOWED_EXTS = {".pdf", ".pptx", ".ppsx", ".doc", ".docx", ".zip", ".png", ".jpg", ".jpeg"}
MAX_SUBMISSION_FILES = 5
MAX_SUBMISSION_FILE_SIZE = 20 * 1024 * 1024
MAX_SUBMISSION_TOTAL_SIZE = 50 * 1024 * 1024


def _parse_due_at(raw: Optional[str]) -> Optional[datetime]:
    if not raw or not raw.strip():
        return None
    text = raw.strip().replace("Z", "+00:00")
    try:
        return datetime.fromisoformat(text)
    except ValueError:
        try:
            return datetime.strptime(raw.strip(), "%Y-%m-%dT%H:%M")
        except ValueError as exc:
            raise HTTPException(status_code=400, detail="截止日期格式无效") from exc


def _serialize_attachment(attachment) -> HomeworkAttachmentResponse:
    return HomeworkAttachmentResponse(
        id=attachment.id,
        filename=attachment.filename,
        file_type=attachment.file_type,
        file_size=attachment.file_size or 0,
    )


def _serialize_submission_attachment(attachment) -> HomeworkAttachmentResponse:
    return HomeworkAttachmentResponse(
        id=attachment.id,
        filename=attachment.filename,
        file_type=attachment.file_type,
        file_size=attachment.file_size or 0,
    )


def _serialize_submission(submission, student_name: str | None = None) -> HomeworkSubmissionResponse:
    attachments = [
        _serialize_submission_attachment(item)
        for item in submission.attachments
    ]
    first = attachments[0] if attachments else None
    return HomeworkSubmissionResponse(
        id=submission.id,
        homework_id=submission.homework_id,
        student_id=submission.student_id,
        student_name=student_name,
        content=submission.content,
        filename=first.filename if first else submission.filename,
        file_type=first.file_type if first else submission.file_type,
        file_size=first.file_size if first else (submission.file_size or 0),
        has_file=bool(attachments or submission.file_path),
        attachments=attachments,
        submitted_at=submission.submitted_at.isoformat() if submission.submitted_at else None,
    )


def _serialize_homework(
    hw,
    attachments: list | None = None,
    submission_count: int = 0,
    my_submission: dict | None = None,
) -> HomeworkResponse:
    serialized_attachments = [
        _serialize_attachment(item) for item in (attachments or [])
    ]
    return HomeworkResponse(
        id=hw.id,
        class_id=hw.class_id,
        title=hw.title,
        description=hw.description,
        due_at=hw.due_at.isoformat() if hw.due_at else None,
        attachment_name=hw.attachment_name,
        has_attachment=bool(serialized_attachments or hw.attachment_path),
        attachments=serialized_attachments,
        created_at=hw.created_at.isoformat() if hw.created_at else "",
        submission_count=submission_count,
        my_submission=my_submission,
    )


def _preview_html(filename: str, blocks: list[tuple[str, str]]) -> str:
    sections = "".join(
        f"<section><h2>{escape(title)}</h2><div>{body}</div></section>"
        for title, body in blocks
    )
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{escape(filename)}</title>
  <style>
    body {{ margin: 0; padding: 24px; font-family: system-ui, sans-serif; background: #f5f7fa; color: #172033; }}
    main {{ width: min(960px, 100%); margin: auto; }}
    h1 {{ font-size: 22px; }}
    section {{ margin: 16px 0; padding: 20px; border-radius: 12px; background: white; box-shadow: 0 4px 18px rgba(20,30,50,.08); }}
    h2 {{ margin-top: 0; font-size: 15px; color: #456; }}
    div {{ line-height: 1.7; white-space: pre-wrap; overflow-wrap: anywhere; }}
  </style>
</head>
<body><main><h1>{escape(filename)}</h1>{sections}</main></body>
</html>"""


def _build_office_preview(file_path: str, filename: str, ext: str) -> str:
    blocks: list[tuple[str, str]] = []
    if ext in {".pptx", ".ppsx"}:
        presentation = Presentation(file_path)
        for index, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            blocks.append((f"幻灯片 {index}", escape(text or "（无可提取文字）")))
    elif ext == ".docx":
        document = Document(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        blocks.append(("文档内容", escape(text or "（无可提取文字）")))
    return _preview_html(filename, blocks)


def _serve_previewable_file(
    file_path: str,
    filename: str,
    download: bool,
):
    ext = os.path.splitext(filename)[1].lower()
    media_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
    if download:
        return FileResponse(
            file_path,
            filename=filename,
            media_type=media_type,
            content_disposition_type="attachment",
        )
    if ext in {".pdf", ".png", ".jpg", ".jpeg"}:
        return FileResponse(
            file_path,
            filename=filename,
            media_type=media_type,
            content_disposition_type="inline",
        )
    if ext in {".pptx", ".ppsx", ".docx"}:
        try:
            return HTMLResponse(_build_office_preview(file_path, filename, ext))
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"生成预览失败: {exc}") from exc
    raise HTTPException(status_code=400, detail="该附件格式不支持在线预览，请下载后查看")


def _remove_file_if_safe(path: str | None) -> None:
    if not path:
        return
    resolved = os.path.realpath(path)
    base_dir = os.path.realpath(settings.HOMEWORK_DIR)
    if os.path.commonpath([resolved, base_dir]) == base_dir and os.path.isfile(resolved):
        os.remove(resolved)


@router.get("/classes/{class_id}/homeworks", response_model=list[HomeworkResponse])
def list_homeworks(
    class_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not class_repo.user_can_access_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="无权访问该班级")

    items = homework_repo.list_homeworks(db, class_id)
    result = []
    for hw in items:
        subs = homework_repo.list_submissions(db, hw.id)
        my_sub = None
        if current_user.role == UserRole.STUDENT:
            mine = homework_repo.get_student_submission(db, hw.id, current_user.id)
            if mine:
                mine_attachments = homework_repo.list_submission_attachments(db, mine.id)
                my_sub = {
                    "id": mine.id,
                    "content": mine.content,
                    "filename": mine_attachments[0].filename if mine_attachments else mine.filename,
                    "has_file": bool(mine_attachments or mine.file_path),
                    "attachments": [
                        _serialize_submission_attachment(item).model_dump()
                        for item in mine_attachments
                    ],
                    "submitted_at": mine.submitted_at.isoformat() if mine.submitted_at else None,
                }
        attachments = homework_repo.list_attachments(db, hw.id)
        result.append(_serialize_homework(
            hw,
            attachments=attachments,
            submission_count=len(subs),
            my_submission=my_sub,
        ))
    return result


@router.post("/classes/{class_id}/homeworks", response_model=HomeworkResponse)
async def create_homework(
    class_id: int,
    title: str = Form(...),
    description: str = Form(""),
    due_at: str = Form(""),
    files: list[UploadFile] | None = File(None),
    file: UploadFile | None = File(None),
    current_user: User = Depends(require_role(UserRole.TEACHER)),
    db: Session = Depends(get_db),
):
    if not class_repo.user_owns_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="仅班级教师可发布作业")
    if not title.strip():
        raise HTTPException(status_code=400, detail="请填写作业标题")

    incoming_files = [item for item in (files or []) if item and item.filename]
    if file and file.filename:
        incoming_files.append(file)

    normalized_files: list[tuple[UploadFile, str, str]] = []
    for upload in incoming_files:
        filename = os.path.basename(upload.filename or "").strip()
        ext = os.path.splitext(filename)[1].lower()
        if ext not in ALLOWED_EXTS:
            raise HTTPException(status_code=400, detail=f"不支持的附件格式: {filename}")
        normalized_files.append((upload, filename, ext))

    hw = homework_repo.create_homework(
        db=db,
        class_id=class_id,
        title=title,
        description=description,
        due_at=_parse_due_at(due_at),
        created_by=current_user.id,
        attachment_path=None,
        attachment_name=None,
    )
    saved_paths: list[str] = []
    attachments = []
    try:
        for upload, filename, ext in normalized_files:
            content = await upload.read()
            saved = data_manager.save_homework_attachment(
                content,
                filename,
                class_id,
                hw.id,
            )
            if saved["status"] != "success":
                raise RuntimeError(saved.get("message", "附件保存失败"))
            saved_paths.append(saved["path"])
            attachments.append(homework_repo.add_attachment(
                db,
                homework_id=hw.id,
                filename=filename,
                file_path=saved["path"],
                file_type=ext.lstrip("."),
                file_size=len(content),
            ))

        if attachments:
            hw.attachment_path = attachments[0].file_path
            hw.attachment_name = attachments[0].filename
            db.commit()
            db.refresh(hw)
    except Exception as exc:
        for path in saved_paths:
            _remove_file_if_safe(path)
        homework_repo.delete_homework(db, hw)
        raise HTTPException(status_code=500, detail=f"发布作业失败: {exc}") from exc

    return _serialize_homework(hw, attachments=attachments)


@router.delete("/homeworks/{homework_id}")
def delete_homework(
    homework_id: int,
    current_user: User = Depends(require_role(UserRole.TEACHER)),
    db: Session = Depends(get_db),
):
    hw = homework_repo.get_homework(db, homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")
    if not class_repo.user_owns_class(db, current_user, hw.class_id):
        raise HTTPException(status_code=403, detail="无权删除该作业")
    paths = {
        hw.attachment_path,
        *(attachment.file_path for attachment in hw.attachments),
        *(submission.file_path for submission in hw.submissions),
        *(
            attachment.file_path
            for submission in hw.submissions
            for attachment in submission.attachments
        ),
    }
    for path in paths:
        _remove_file_if_safe(path)
    assignment_dir = os.path.join(
        settings.HOMEWORK_DIR,
        str(hw.class_id),
        "assignments",
        str(hw.id),
    )
    resolved_dir = os.path.realpath(assignment_dir)
    base_dir = os.path.realpath(settings.HOMEWORK_DIR)
    if os.path.commonpath([resolved_dir, base_dir]) == base_dir and os.path.isdir(resolved_dir):
        shutil.rmtree(resolved_dir)
    homework_repo.delete_homework(db, hw)
    return {"ok": True}


@router.get("/homeworks/{homework_id}/attachments/{attachment_id}/file")
def get_homework_attachment_file(
    homework_id: int,
    attachment_id: int,
    download: bool = Query(False),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    hw = homework_repo.get_homework(db, homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")
    if not class_repo.user_can_access_class(db, current_user, hw.class_id):
        raise HTTPException(status_code=403, detail="无权访问")
    attachment = homework_repo.get_attachment(db, homework_id, attachment_id)
    if not attachment:
        raise HTTPException(status_code=404, detail="附件不存在")

    file_path = os.path.realpath(attachment.file_path)
    base_dir = os.path.realpath(settings.HOMEWORK_DIR)
    if os.path.commonpath([file_path, base_dir]) != base_dir or not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="附件文件缺失")
    ext = os.path.splitext(attachment.filename)[1].lower()
    media_type = mimetypes.guess_type(attachment.filename)[0] or "application/octet-stream"

    if download:
        return FileResponse(
            file_path,
            filename=attachment.filename,
            media_type=media_type,
            content_disposition_type="attachment",
        )
    if ext in {".pdf", ".png", ".jpg", ".jpeg"}:
        return FileResponse(
            file_path,
            filename=attachment.filename,
            media_type=media_type,
            content_disposition_type="inline",
        )
    if ext in {".pptx", ".ppsx", ".docx"}:
        try:
            return HTMLResponse(_build_office_preview(
                file_path, attachment.filename, ext
            ))
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"生成预览失败: {exc}") from exc
    raise HTTPException(status_code=400, detail="该附件格式不支持在线预览，请下载后查看")


@router.get("/homeworks/{homework_id}/attachment")
def download_homework_attachment(
    homework_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    hw = homework_repo.get_homework(db, homework_id)
    if not hw or not hw.attachment_path:
        raise HTTPException(status_code=404, detail="附件不存在")
    if not class_repo.user_can_access_class(db, current_user, hw.class_id):
        raise HTTPException(status_code=403, detail="无权访问")
    if not os.path.isfile(hw.attachment_path):
        raise HTTPException(status_code=404, detail="附件文件缺失")
    return FileResponse(
        hw.attachment_path,
        filename=hw.attachment_name or os.path.basename(hw.attachment_path),
        media_type=mimetypes.guess_type(hw.attachment_path)[0] or "application/octet-stream",
    )


@router.post("/homeworks/{homework_id}/submit", response_model=HomeworkSubmissionResponse)
async def submit_homework(
    homework_id: int,
    content: str = Form(""),
    files: list[UploadFile] | None = File(None),
    file: UploadFile | None = File(None),
    retained_attachment_ids: str | None = Form(None),
    current_user: User = Depends(require_role(UserRole.STUDENT)),
    db: Session = Depends(get_db),
):
    hw = homework_repo.get_homework(db, homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")
    if not class_repo.user_can_access_class(db, current_user, hw.class_id):
        raise HTTPException(status_code=403, detail="无权提交该作业")

    existing = homework_repo.get_student_submission(db, homework_id, current_user.id)
    existing_attachments = (
        homework_repo.list_submission_attachments(db, existing.id)
        if existing else []
    )
    existing_by_id = {item.id: item for item in existing_attachments}

    if retained_attachment_ids is None:
        retained_ids = set(existing_by_id)
    else:
        try:
            parsed_ids = json.loads(retained_attachment_ids)
            if not isinstance(parsed_ids, list):
                raise ValueError
            retained_ids = {int(item) for item in parsed_ids}
        except (TypeError, ValueError, json.JSONDecodeError) as exc:
            raise HTTPException(status_code=400, detail="保留附件清单格式无效") from exc
        if not retained_ids.issubset(existing_by_id):
            raise HTTPException(status_code=400, detail="保留附件清单包含无效附件")

    incoming_files = [item for item in (files or []) if item and item.filename]
    if file and file.filename:
        if files or retained_attachment_ids is not None:
            incoming_files.append(file)
        else:
            retained_ids = set()
            incoming_files = [file]

    if len(retained_ids) + len(incoming_files) > MAX_SUBMISSION_FILES:
        raise HTTPException(status_code=400, detail="每份作业最多提交5个附件")

    prepared_files: list[dict] = []
    retained_total = sum(existing_by_id[item].file_size or 0 for item in retained_ids)
    new_total = 0
    for upload in incoming_files:
        filename = os.path.basename(upload.filename or "").strip()
        ext = os.path.splitext(filename)[1].lower()
        if not filename or ext not in ALLOWED_EXTS:
            raise HTTPException(status_code=400, detail=f"不支持的提交文件格式: {filename}")
        raw = await upload.read()
        if len(raw) > MAX_SUBMISSION_FILE_SIZE:
            raise HTTPException(status_code=400, detail=f"附件超过20MB限制: {filename}")
        new_total += len(raw)
        if retained_total + new_total > MAX_SUBMISSION_TOTAL_SIZE:
            raise HTTPException(status_code=400, detail="全部附件合计不能超过50MB")
        prepared_files.append({
            "filename": filename,
            "file_type": ext.lstrip("."),
            "file_size": len(raw),
            "content": raw,
        })

    if not (content or "").strip() and not retained_ids and not prepared_files:
        raise HTTPException(status_code=400, detail="请填写文字说明或上传文件")

    saved_paths: list[str] = []
    new_attachments: list[dict] = []
    try:
        for item in prepared_files:
            saved = data_manager.save_submission_attachment(
                item["content"],
                item["filename"],
                hw.class_id,
                hw.id,
                current_user.id,
            )
            if saved["status"] != "success":
                raise RuntimeError(saved.get("message", "提交文件保存失败"))
            saved_paths.append(saved["path"])
            new_attachments.append({
                **item,
                "file_path": saved["path"],
            })

        sub, removed_paths = homework_repo.save_submission_with_attachments(
            db=db,
            homework_id=homework_id,
            student_id=current_user.id,
            content=content,
            retained_attachment_ids=retained_ids,
            new_attachments=new_attachments,
        )
    except HTTPException:
        for path in saved_paths:
            _remove_file_if_safe(path)
        raise
    except Exception as exc:
        for path in saved_paths:
            _remove_file_if_safe(path)
        raise HTTPException(status_code=500, detail=f"提交作业失败: {exc}") from exc

    for path in removed_paths:
        _remove_file_if_safe(path)
    db.refresh(sub)
    return _serialize_submission(sub, current_user.name)


@router.get("/homeworks/{homework_id}/submissions", response_model=list[HomeworkSubmissionResponse])
def list_homework_submissions(
    homework_id: int,
    current_user: User = Depends(require_role(UserRole.TEACHER)),
    db: Session = Depends(get_db),
):
    hw = homework_repo.get_homework(db, homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")
    if not class_repo.user_owns_class(db, current_user, hw.class_id):
        raise HTTPException(status_code=403, detail="无权查看提交")
    rows = homework_repo.list_submissions(db, homework_id)
    return [HomeworkSubmissionResponse(**row) for row in rows]


@router.get("/submissions/{submission_id}/attachments/{attachment_id}/file")
def get_submission_attachment_file(
    submission_id: int,
    attachment_id: int,
    download: bool = Query(False),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    sub = homework_repo.get_submission(db, submission_id)
    if not sub:
        raise HTTPException(status_code=404, detail="提交不存在")
    hw = homework_repo.get_homework(db, sub.homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")

    is_owner = class_repo.user_owns_class(db, current_user, hw.class_id)
    is_self = current_user.id == sub.student_id
    if not (is_owner or is_self):
        raise HTTPException(status_code=403, detail="无权访问该提交附件")

    attachment = homework_repo.get_submission_attachment(
        db,
        submission_id,
        attachment_id,
    )
    if not attachment:
        raise HTTPException(status_code=404, detail="提交附件不存在")

    file_path = os.path.realpath(attachment.file_path)
    base_dir = os.path.realpath(settings.HOMEWORK_DIR)
    if os.path.commonpath([file_path, base_dir]) != base_dir or not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="提交附件文件缺失")
    return _serve_previewable_file(file_path, attachment.filename, download)


@router.get("/submissions/{submission_id}/file")
def download_submission_file(
    submission_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    sub = homework_repo.get_submission(db, submission_id)
    if not sub or not sub.file_path:
        raise HTTPException(status_code=404, detail="提交文件不存在")
    hw = homework_repo.get_homework(db, sub.homework_id)
    if not hw:
        raise HTTPException(status_code=404, detail="作业不存在")

    is_owner = class_repo.user_owns_class(db, current_user, hw.class_id)
    is_self = current_user.id == sub.student_id
    if not (is_owner or is_self):
        raise HTTPException(status_code=403, detail="无权下载该提交")
    file_path = os.path.realpath(sub.file_path)
    base_dir = os.path.realpath(settings.HOMEWORK_DIR)
    if os.path.commonpath([file_path, base_dir]) != base_dir or not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="文件缺失")

    return FileResponse(
        file_path,
        filename=sub.filename or os.path.basename(file_path),
        media_type=mimetypes.guess_type(file_path)[0] or "application/octet-stream",
        content_disposition_type="attachment",
    )
