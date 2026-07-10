import os
from html import escape
from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Query
from fastapi.responses import FileResponse, HTMLResponse
from sqlalchemy.orm import Session
from pptx import Presentation
from agent_core.config.settings import settings
from database.mysql_db import get_db, User, UserRole
from database import class_repo
from app.core.deps import get_current_user, require_role
from .schemas import CreateClassRequest, JoinClassRequest, ClassResponse, MaterialResponse

router = APIRouter()

MIME_BY_EXT = {
    ".pdf": "application/pdf",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppsx": "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
}


def _build_pptx_preview_html(file_path: str, filename: str) -> str:
    prs = Presentation(file_path)
    slide_blocks = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(escape(text))
        body = "<br>".join(line.replace("\n", "<br>") for line in lines)
        if not body:
            body = "<span class='empty'>No text content extracted from this slide.</span>"
        slide_blocks.append(f"""
          <section class="slide">
            <div class="slide-index">Slide {index}</div>
            <div class="slide-body">{body}</div>
          </section>
        """)

    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{escape(filename)} preview</title>
  <style>
    body {{ margin: 0; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #f4f8fb; color: #16304a; }}
    header {{ position: sticky; top: 0; padding: 18px 28px; background: rgba(244, 248, 251, 0.92); border-bottom: 1px solid rgba(22, 48, 74, 0.12); backdrop-filter: blur(12px); }}
    h1 {{ margin: 0; font-size: 20px; }}
    main {{ width: min(960px, calc(100% - 32px)); margin: 24px auto 48px; display: grid; gap: 18px; }}
    .slide {{ background: #fff; border: 1px solid rgba(22, 48, 74, 0.12); border-radius: 8px; padding: 22px 24px; box-shadow: 0 12px 32px rgba(22, 48, 74, 0.08); }}
    .slide-index {{ font-size: 12px; font-weight: 700; letter-spacing: 0.08em; text-transform: uppercase; color: #4a86b8; margin-bottom: 12px; }}
    .slide-body {{ line-height: 1.7; white-space: normal; }}
    .empty {{ color: #7d8fa3; }}
  </style>
</head>
<body>
  <header><h1>{escape(filename)}</h1></header>
  <main>{''.join(slide_blocks)}</main>
</body>
</html>"""


@router.post("", response_model=ClassResponse)
async def create_class(
    payload: CreateClassRequest,
    current_user: User = Depends(require_role(UserRole.TEACHER)),
    db: Session = Depends(get_db),
):
    if not payload.name.strip():
        raise HTTPException(status_code=400, detail="班级名称不能为空")
    classroom = class_repo.create_class(db, current_user, payload.name)
    return ClassResponse(
        id=classroom.id,
        name=classroom.name,
        invite_code=classroom.invite_code,
        teacher_id=classroom.teacher_id,
        role_in_class="owner",
    )


@router.get("/mine", response_model=list[ClassResponse])
async def list_my_classes(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if current_user.role is None:
        raise HTTPException(status_code=400, detail="请先选择学生或教师身份")
    return class_repo.get_user_classes(db, current_user)


@router.post("/join", response_model=ClassResponse)
async def join_class(
    payload: JoinClassRequest,
    current_user: User = Depends(require_role(UserRole.STUDENT)),
    db: Session = Depends(get_db),
):
    try:
        classroom = class_repo.join_class(db, current_user, payload.invite_code)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    return ClassResponse(
        id=classroom.id,
        name=classroom.name,
        invite_code=classroom.invite_code,
        teacher_id=classroom.teacher_id,
        role_in_class="member",
    )


@router.get("/{class_id}/materials", response_model=list[MaterialResponse])
async def list_class_materials(
    class_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not class_repo.user_can_access_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="无权访问该班级资料")
    materials = class_repo.list_materials(db, class_id)
    return [
        MaterialResponse(
            id=m.id,
            class_id=m.class_id,
            filename=m.filename,
            file_type=m.file_type,
            file_size=m.file_size,
            uploaded_at=m.uploaded_at.isoformat(),
        )
        for m in materials
    ]


@router.post("/{class_id}/materials")
async def upload_class_material(
    class_id: int,
    file: UploadFile = File(...),
    current_user: User = Depends(require_role(UserRole.TEACHER)),
    db: Session = Depends(get_db),
):
    from app.core.data_manager import data_manager

    if not class_repo.user_owns_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="仅班级创建者可上传资料")

    ext = (file.filename or "").lower()
    if not ext.endswith((".pdf", ".pptx", ".ppsx")):
        raise HTTPException(status_code=400, detail="仅支持 PDF、PPTX、PPSX 文件")

    content = await file.read()
    result = data_manager.save_class_material(content, file.filename, class_id)
    if result["status"] != "success":
        raise HTTPException(status_code=500, detail=result.get("message", "上传失败"))

    file_type = "pdf" if ext.endswith(".pdf") else "pptx"
    material = class_repo.add_material(
        db,
        class_id=class_id,
        filename=file.filename,
        file_path=result["path"],
        file_type=file_type,
        file_size=len(content),
        uploader_id=current_user.id,
    )
    return {
        "message": f"文件 {file.filename} 上传成功",
        "material": MaterialResponse(
            id=material.id,
            class_id=material.class_id,
            filename=material.filename,
            file_type=material.file_type,
            file_size=material.file_size,
            uploaded_at=material.uploaded_at.isoformat(),
        ),
    }


@router.get("/{class_id}/materials/{material_id}/file")
async def get_class_material_file(
    class_id: int,
    material_id: int,
    download: bool = Query(False, description="true 为下载，false 为在线预览（PDF）"),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not class_repo.user_can_access_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="无权访问该班级资料")

    material = class_repo.get_material(db, class_id, material_id)
    if not material:
        raise HTTPException(status_code=404, detail="资料不存在")

    file_path = os.path.realpath(material.file_path)
    base_dir = os.path.realpath(settings.CLASS_MATERIALS_DIR)
    if not file_path.startswith(base_dir) or not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="文件不存在或已损坏")

    ext = os.path.splitext(material.filename)[1].lower()
    media_type = MIME_BY_EXT.get(ext, "application/octet-stream")
    disposition = "attachment" if download or ext != ".pdf" else "inline"

    return FileResponse(
        path=file_path,
        media_type=media_type,
        filename=material.filename,
        content_disposition_type=disposition,
    )


@router.get("/{class_id}/materials/{material_id}/preview", response_class=HTMLResponse)
async def preview_class_material(
    class_id: int,
    material_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not class_repo.user_can_access_class(db, current_user, class_id):
        raise HTTPException(status_code=403, detail="无权访问该班级资料")

    material = class_repo.get_material(db, class_id, material_id)
    if not material:
        raise HTTPException(status_code=404, detail="资料不存在")

    file_path = os.path.realpath(material.file_path)
    base_dir = os.path.realpath(settings.CLASS_MATERIALS_DIR)
    if not file_path.startswith(base_dir) or not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="文件不存在或已损坏")

    ext = os.path.splitext(material.filename)[1].lower()
    if ext == ".pdf":
        return FileResponse(
            path=file_path,
            media_type="application/pdf",
            filename=material.filename,
            content_disposition_type="inline",
        )
    if ext in {".pptx", ".ppsx"}:
        try:
            return HTMLResponse(_build_pptx_preview_html(file_path, material.filename))
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"生成预览失败: {e}")

    raise HTTPException(status_code=400, detail="该文件类型暂不支持在线预览")
