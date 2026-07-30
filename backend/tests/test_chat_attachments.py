import importlib.util
from io import BytesIO
from pathlib import Path
import tempfile
from types import SimpleNamespace
import unittest
from unittest.mock import MagicMock, patch
import zipfile

from alembic.migration import MigrationContext
from alembic.operations import Operations
import fitz
from docx import Document
from fastapi import HTTPException
from pptx import Presentation
from sqlalchemy import create_engine, inspect, text
from sqlalchemy.orm import sessionmaker

from agent_core.config.settings import settings
from app.interfaces.api import chat_attachment_routes
from app.services import chat_attachment_context, chat_attachment_service
from database import chat_attachment_repo
from database.mysql_db import (
    Base,
    ChatMessage,
    ChatMessageAttachment,
    Conversation,
    User,
    UserRole,
)


class ChatAttachmentServiceTests(unittest.TestCase):
    def setUp(self):
        self.engine = create_engine("sqlite:///:memory:")
        Base.metadata.create_all(self.engine)
        self.Session = sessionmaker(bind=self.engine)
        self.db = self.Session()
        self.user = User(
            email="attachment-owner@tongji.edu.cn",
            name="Owner",
            hashed_password="x",
            role=UserRole.STUDENT,
        )
        self.other_user = User(
            email="other-user@tongji.edu.cn",
            name="Other",
            hashed_password="x",
            role=UserRole.STUDENT,
        )
        self.db.add_all([self.user, self.other_user])
        self.db.flush()
        self.conversation = Conversation(
            user_id=self.user.id,
            title="Documents",
        )
        self.other_conversation = Conversation(
            user_id=self.other_user.id,
            title="Private",
        )
        self.db.add_all([self.conversation, self.other_conversation])
        self.db.commit()
        self.temp_dir = tempfile.TemporaryDirectory()
        self.settings_patch = patch.object(
            settings,
            "CHAT_ATTACHMENTS_DIR",
            self.temp_dir.name,
        )
        self.settings_patch.start()

    def tearDown(self):
        self.settings_patch.stop()
        self.temp_dir.cleanup()
        self.db.close()
        self.engine.dispose()

    def _attachment(self, filename, file_type, raw):
        public_id = f"00000000-0000-4000-8000-{self.db.query(ChatMessageAttachment).count() + 1:012d}"
        path = chat_attachment_service.save_attachment_file(
            self.user.id,
            public_id,
            f".{file_type}",
            raw,
        )
        return chat_attachment_repo.create_attachment(
            self.db,
            public_id=public_id,
            user_id=self.user.id,
            filename=filename,
            file_path=path,
            file_type=file_type,
            mime_type=chat_attachment_service.DOCUMENT_MIME_TYPES[f".{file_type}"],
            file_size=len(raw),
        )

    def test_validates_modern_formats_and_rejects_legacy_or_spoofed_files(self):
        self.assertEqual(
            chat_attachment_service.validate_document_bytes(
                "notes.pdf",
                b"%PDF-1.7\n",
            )[0],
            ".pdf",
        )
        office_buffer = BytesIO()
        with zipfile.ZipFile(office_buffer, "w") as archive:
            archive.writestr("[Content_Types].xml", "<Types />")
        self.assertEqual(
            chat_attachment_service.validate_document_bytes(
                "slides.pptx",
                office_buffer.getvalue(),
            )[0],
            ".pptx",
        )
        for filename in ("legacy.doc", "legacy.ppt", "notes.txt"):
            with self.subTest(filename=filename), self.assertRaises(
                chat_attachment_service.AttachmentParseError
            ):
                chat_attachment_service.validate_document_bytes(filename, b"data")
        with self.assertRaises(chat_attachment_service.AttachmentParseError):
            chat_attachment_service.validate_document_bytes("fake.pdf", b"not a pdf")

    def test_extracts_docx_paragraphs_and_tables(self):
        buffer = BytesIO()
        document = Document()
        document.add_paragraph("离散数学课程说明")
        table = document.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "概念"
        table.cell(0, 1).text = "定义"
        document.save(buffer)
        attachment = self._attachment("notes.docx", "docx", buffer.getvalue())
        chat_attachment_repo.claim_attachment(self.db, attachment.public_id)

        chunks, truncated = chat_attachment_service._extract_docx(
            self.db,
            attachment,
            chat_attachment_service.resolve_attachment_path(attachment.file_path),
        )

        self.assertFalse(truncated)
        self.assertIn("离散数学课程说明", str(chunks))
        self.assertIn("概念 | 定义", str(chunks))

    def test_extracts_presentation_by_slide(self):
        buffer = BytesIO()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "命题逻辑"
        slide.placeholders[1].text = "合取与析取"
        presentation.save(buffer)
        attachment = self._attachment("logic.pptx", "pptx", buffer.getvalue())
        chat_attachment_repo.claim_attachment(self.db, attachment.public_id)

        chunks, _ = chat_attachment_service._extract_presentation(
            self.db,
            attachment,
            chat_attachment_service.resolve_attachment_path(attachment.file_path),
        )

        self.assertEqual(chunks[0]["location"], "第 1 张幻灯片")
        self.assertIn("命题逻辑", chunks[0]["text"])

    def test_mixed_pdf_only_ocr_missing_text_pages(self):
        document = fitz.open()
        first = document.new_page()
        first.insert_text((72, 72), "This page already has enough native text for extraction.")
        document.new_page()
        raw = document.tobytes()
        document.close()
        attachment = self._attachment("mixed.pdf", "pdf", raw)
        chat_attachment_repo.claim_attachment(self.db, attachment.public_id)

        with patch.object(
            chat_attachment_service,
            "_ocr_pdf_page",
            return_value="第二页扫描内容",
        ) as ocr:
            chunks, _ = chat_attachment_service._extract_pdf(
                self.db,
                attachment,
                chat_attachment_service.resolve_attachment_path(attachment.file_path),
            )

        self.assertEqual(ocr.call_count, 1)
        self.assertEqual(ocr.call_args.args[2], 2)
        self.assertIn("第二页扫描内容", str(chunks))

    def test_scanned_pdf_over_page_limit_is_rejected_before_ocr(self):
        document = fitz.open()
        for _ in range(settings.CHAT_ATTACHMENT_OCR_MAX_PAGES + 1):
            document.new_page()
        raw = document.tobytes()
        document.close()
        attachment = self._attachment("scan.pdf", "pdf", raw)
        chat_attachment_repo.claim_attachment(self.db, attachment.public_id)

        with (
            patch.object(chat_attachment_service, "_ocr_pdf_page") as ocr,
            self.assertRaisesRegex(
                chat_attachment_service.AttachmentParseError,
                "请拆分后上传",
            ),
        ):
            chat_attachment_service._extract_pdf(
                self.db,
                attachment,
                chat_attachment_service.resolve_attachment_path(attachment.file_path),
            )
        ocr.assert_not_called()

    def test_ocr_retries_twice_then_returns_text(self):
        document = fitz.open()
        page = document.new_page()
        completions = MagicMock()
        completions.create.side_effect = [
            RuntimeError("temporary 1"),
            RuntimeError("temporary 2"),
            SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(content="第三次成功")
                    )
                ]
            ),
        ]
        client = SimpleNamespace(
            chat=SimpleNamespace(completions=completions)
        )
        with (
            patch.object(chat_attachment_service, "OpenAI", return_value=client),
            patch.object(chat_attachment_service.time, "sleep"),
        ):
            result = chat_attachment_service._ocr_pdf_page(
                page,
                "scan.pdf",
                1,
            )
        document.close()

        self.assertEqual(result, "第三次成功")
        self.assertEqual(completions.create.call_count, 3)

    def test_context_search_is_scoped_to_current_conversation(self):
        message = ChatMessage(
            conversation_id=self.conversation.id,
            role="user",
            content="[文档]",
        )
        private_message = ChatMessage(
            conversation_id=self.other_conversation.id,
            role="user",
            content="[文档]",
        )
        self.db.add_all([message, private_message])
        self.db.flush()
        self.db.add_all([
            ChatMessageAttachment(
                public_id="10000000-0000-4000-8000-000000000001",
                user_id=self.user.id,
                message_id=message.id,
                filename="graph.pdf",
                file_path="1/graph.pdf",
                file_type="pdf",
                mime_type="application/pdf",
                file_size=10,
                status="ready",
                extracted_content=[{
                    "location": "第 2 页",
                    "text": "欧拉回路经过图中的每一条边一次。",
                }],
            ),
            ChatMessageAttachment(
                public_id="20000000-0000-4000-8000-000000000001",
                user_id=self.other_user.id,
                message_id=private_message.id,
                filename="private.pdf",
                file_path="2/private.pdf",
                file_type="pdf",
                mime_type="application/pdf",
                file_size=10,
                status="ready",
                extracted_content=[{
                    "location": "第 1 页",
                    "text": "不能出现在其他用户上下文中的秘密。",
                }],
            ),
        ])
        self.db.commit()

        reference = chat_attachment_context.build_document_reference(
            self.db,
            conversation_id=self.conversation.id,
            question="什么是欧拉回路？",
            current_message_id=message.id,
        )

        self.assertIn("graph.pdf | 第 2 页", reference)
        self.assertIn("欧拉回路", reference)
        self.assertNotIn("秘密", reference)
        self.assertIn("未经信任的参考资料", reference)

    def test_owner_filter_and_cancel_prevent_cross_user_access(self):
        attachment = self._attachment("notes.pdf", "pdf", b"%PDF-1.7\n")
        self.assertIsNone(
            chat_attachment_repo.get_owned_attachment(
                self.db,
                attachment.public_id,
                self.other_user.id,
            )
        )
        cancelled = chat_attachment_repo.cancel_unassociated(
            self.db,
            attachment.public_id,
            self.user.id,
        )
        self.assertEqual(cancelled.status, "cancelled")
        self.assertTrue(
            chat_attachment_repo.is_cancelled(self.db, attachment.public_id)
        )

    def test_private_file_route_rejects_another_user_and_path_traversal(self):
        attachment = self._attachment("notes.pdf", "pdf", b"%PDF-1.7\n")
        with self.assertRaises(HTTPException) as forbidden:
            chat_attachment_routes.get_chat_attachment_file(
                public_id=attachment.public_id,
                download=False,
                current_user=self.other_user,
                db=self.db,
            )
        self.assertEqual(forbidden.exception.status_code, 404)
        with self.assertRaises(chat_attachment_service.AttachmentParseError):
            chat_attachment_service.resolve_attachment_path("../secret.txt")


class ChatAttachmentMigrationTests(unittest.TestCase):
    def test_upgrade_creates_expected_table_and_indexes_idempotently(self):
        migration_path = (
            Path(__file__).resolve().parents[1]
            / "alembic"
            / "versions"
            / "20260730_0005_chat_message_attachments.py"
        )
        spec = importlib.util.spec_from_file_location(
            "chat_attachment_migration",
            migration_path,
        )
        migration = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(migration)
        engine = create_engine("sqlite:///:memory:")
        with engine.begin() as connection:
            connection.execute(text("CREATE TABLE users (id INTEGER PRIMARY KEY)"))
            connection.execute(text(
                "CREATE TABLE chat_messages (id INTEGER PRIMARY KEY)"
            ))
            migration.op = Operations(MigrationContext.configure(connection))
            migration.upgrade()
            migration.upgrade()
            columns = {
                item["name"]
                for item in inspect(connection).get_columns(
                    "chat_message_attachments"
                )
            }
            indexes = {
                item["name"]
                for item in inspect(connection).get_indexes(
                    "chat_message_attachments"
                )
            }
        engine.dispose()

        self.assertIn("extracted_content", columns)
        self.assertIn("message_id", columns)
        self.assertIn("ix_chat_message_attachments_user_id", indexes)
        self.assertIn("ix_chat_message_attachments_status", indexes)


if __name__ == "__main__":
    unittest.main()
